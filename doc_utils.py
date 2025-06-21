from pptx import Presentation
from python_docx_replace import docx_replace
from docx import Document


def replace_text_pptx(replacements, template_path, output_path):
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in replacements.items():
                            placeholder = f"$[{key}]"
                            if placeholder in run.text:
                                run.text = run.text.replace(
                                    placeholder, str(value))
    prs.save(output_path)
    return output_path


def export_agreement(replacements, template_path, output_path):
    doc = Document(template_path)
    docx_replace(doc, **replacements)
    doc.save(output_path)
